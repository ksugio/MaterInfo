from django.db import models
from django.urls import reverse
from project.models import Project, Created, Updated, Remote, UpperModelUploadTo
from io import StringIO
from zipfile import ZipFile
from nbconvert import HTMLExporter
import os
import docx
import pptx
import pandas as pd

def DocumentUploadTo(instance, filename):
    return filename

class Document(Created, Updated, Remote):
    upper = models.ForeignKey(Project, verbose_name='Project', on_delete=models.CASCADE)
    title = models.CharField(verbose_name='Title', max_length=100)
    StatusChoices = ((0, 'Active'), (1, 'Stop'), (2, 'Finish'))
    status = models.PositiveSmallIntegerField(verbose_name='Status', choices=StatusChoices, default=0)
    note = models.TextField(verbose_name='Note', blank=True)

    def __str__(self):
        return self.title

    def get_list_url(self):
        return reverse('document:list', kwargs={'pk': self.upper.id})

    def get_detail_url(self):
        return reverse('document:detail', kwargs={'pk': self.id, 'ed': self.latest_file().edition, 'zipid': 0})

    def get_update_url(self):
        return reverse('document:update', kwargs={'pk': self.id})

    def get_delete_url(self):
        return reverse('document:delete', kwargs={'pk': self.id})

    def get_apiupdate_url(self):
        return reverse('document:api_update', kwargs={'pk': self.id})

    def latest_file(self):
        files = File.objects.filter(upper=self).order_by('-edition')
        if files:
            return files[0]
        else:
            return None

class File(Created, Remote):
    upper = models.ForeignKey(Document, verbose_name='Document', on_delete=models.CASCADE)
    file = models.FileField(verbose_name='File', upload_to=UpperModelUploadTo)
    comment = models.TextField(verbose_name='Comment', blank=True)
    edition = models.PositiveSmallIntegerField(verbose_name='Edition')
    filename = models.CharField(verbose_name='Filename', max_length=256)
    CodeExtentions = ['.c', '.cpp', '.cs', '.css', '.html', '.htm',
                      '.java', '.js', '.sql', '.pl', '.php',
                      '.py', '.r', '.rb', '.rs', '.sh', '.ts',
                      '.xml', '.yaml']
    ImageExtentions = {'.jpg': 'image/jpeg', '.jpeg': 'image/jpeg',
                       '.JPG': 'image/jpeg', '.JPEG': 'image/jpeg',
                       '.png': 'image/png', '.PNG': 'image/png'}

    def get_detail_url(self):
        return reverse('document:update', kwargs={'pk': self.upper.id})

    def get_update_url(self):
        return reverse('document:file_update', kwargs={'pk': self.id})

    def basename(self):
        return os.path.basename(self.file.name)

    def title(self):
        return '{} : Edition {} Comment'.format(self.upper.title, self.edition)

    def file_ext(self):
        fn, ext = os.path.splitext(self.file.name)
        return ext

    def file_content(self):
        fn, ext = os.path.splitext(self.file.name)
        with self.file.open('rb') as f:
            content = self.content(f, ext)
        content['zipid'] = 0
        return content

    def zipfile_list(self):
        with self.file.open('rb') as f:
            zf = ZipFile(f)
            return zf.namelist()

    def zipfile_content(self, zipid):
        fn, ext = os.path.splitext(self.file.name)
        if ext == '.zip':
            with self.file.open('rb') as f:
                zf = ZipFile(f)
                names = zf.namelist()
                if zipid > 0 and zipid <= len(names):
                    name = names[zipid - 1]
                    fn, ext = os.path.splitext(name)
                    with zf.open(name) as ff:
                        content = self.content(ff, ext)
                    content['zipname'] = name
                    content['listnum'] = len(names)
                    content['zipid'] = zipid
                    return content

    def content(self, f, ext):
        if ext == '.docx':
            return {
                'ext': ext,
                'texts': self.docx_content(f)
            }
        elif ext == '.pptx':
            return {
                'ext': ext,
                'texts': self.pptx_content(f)
            }
        elif ext == '.xlsx':
            return {
                'ext': ext,
                'html': self.xlsx_content(f)
            }
        elif ext == '.csv':
            return {
                'ext': ext,
                'html': self.csv_content(f)
            }
        elif ext == '.md':
            return {
                'ext': ext,
                'md': f.read().decode('utf-8')
            }
        elif ext in self.CodeExtentions:
            return {
                'ext': ext,
                'code': f.read().decode('utf-8')
            }
        elif ext == '.ipynb':
            body, _ = HTMLExporter().from_file(f)
            return {
                'ext': ext,
                'html': body
            }
        elif ext in self.ImageExtentions:
            return {
                'ext': ext,
                'type': self.ImageExtentions[ext],
                'image': f.read(),
                'filename': self.filename
            }
        elif ext == '.pdf':
            return {
                'ext': ext,
                'type': 'application/pdf',
                'image': f.read(),
                'filename': self.filename
            }
        else:
            dat = f.read()
            if b"\x00" not in dat[:516]:
                return {
                    'ext': ext,
                    'plain': dat.decode('utf-8')
                }
            else:
                return {
                    'ext': ext
                }

    def is_bin(self, f):
        ret = f.read()



    def docx_content(self, f):
        doc = docx.Document(f)
        texts = []
        for para in doc.paragraphs:
            texts.append(para.text)
        return texts

    def pptx_content(self, f):
        prs = pptx.Presentation(f)
        texts = []
        for i, sld in enumerate(prs.slides):
            texts.append('---------- Page {0} ----------'.format(i + 1))
            for shape in sld.shapes:
                if shape.has_text_frame:
                    for text in shape.text.splitlines():
                        texts.append(text)
                if shape.has_table:
                    for cell in shape.table.iter_cells():
                        for text in cell.text.splitlines():
                            texts.append(text)
        return texts

    def xlsx_content(self, f):
        df = pd.read_excel(f)
        buf = StringIO()
        df.to_html(buf)
        table = buf.getvalue()
        buf.close()
        return table

    def csv_content(self, f):
        df = pd.read_csv(f)
        buf = StringIO()
        df.to_html(buf)
        table = buf.getvalue()
        buf.close()
        return table

